from __future__ import annotations

import json
import os
import re
import shutil
from pathlib import Path

from _bootstrap import ROOT

os.environ.setdefault("MPLBACKEND", "Agg")
os.environ.setdefault("MPLCONFIGDIR", str(ROOT / "artifacts" / ".matplotlib"))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt

from build_report_from_template import build_report
from office_hygiene import scrub_office_metadata

EVAL_DIR = ROOT / "artifacts" / "evaluation"
SUBMISSION_DIR = ROOT / "submission"

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
plt.rcParams["axes.unicode_minus"] = False


def load_inputs():
    metrics = json.loads((EVAL_DIR / "metrics.json").read_text(encoding="utf-8"))
    training = json.loads((ROOT / "artifacts" / "training_result.json").read_text(encoding="utf-8"))
    alerts = []
    alerts_path = EVAL_DIR / "alerts.jsonl"
    if alerts_path.exists():
        for line in alerts_path.read_text(encoding="utf-8").splitlines()[:20]:
            if line.strip():
                alerts.append(json.loads(line))
    return metrics, training, alerts


def make_figures(metrics, training):
    EVAL_DIR.mkdir(parents=True, exist_ok=True)
    colors = {"navy": "#0d1b2a", "blue": "#4f8bd6", "cyan": "#42c9c4", "red": "#e45b5b", "green": "#39a77a"}
    fig, ax = plt.subplots(figsize=(12, 4.8))
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.8); ax.axis("off")
    stages = [
        (0.25, "日志接入", "Agent + 主机事件"),
        (2.7, "序列编码", "行为词元 + 时间特征"),
        (5.15, "轻量 Transformer", "注意力 + 下一事件"),
        (7.6, "误报率校准", "异常分数 + 告警阈值"),
        (10.05, "证据告警", "自然语言解释 + 原始日志"),
    ]
    for index, (x, title, subtitle) in enumerate(stages):
        box = FancyBboxPatch((x, 1.55), 1.8, 1.65, boxstyle="round,pad=0.08,rounding_size=0.12", fc="#eef5fc", ec=colors["blue"], lw=1.8)
        ax.add_patch(box); ax.text(x+0.9, 2.63, title, ha="center", va="center", fontsize=11, weight="bold", color=colors["navy"]); ax.text(x+0.9, 2.08, subtitle, ha="center", va="center", fontsize=9, color="#4e647d")
        if index < len(stages)-1:
            ax.add_patch(FancyArrowPatch((x+1.84, 2.36), (x+2.38, 2.36), arrowstyle="-|>", mutation_scale=15, color=colors["cyan"], lw=2))
    ax.text(6, 4.25, "AgentGuard 端到端系统架构", ha="center", fontsize=17, weight="bold", color=colors["navy"])
    ax.text(6, .62, "离线运行 · 可复现 · CPU部署 · 证据可追溯", ha="center", fontsize=11, color="#60758c")
    fig.tight_layout(); fig.savefig(EVAL_DIR / "model_architecture.png", dpi=190, bbox_inches="tight"); plt.close(fig)

    matrix = [[metrics["true_negative"], metrics["false_positive"]], [metrics["false_negative"], metrics["true_positive"]]]
    fig, ax = plt.subplots(figsize=(5.8, 5.2)); image = ax.imshow(matrix, cmap="Blues")
    ax.set_xticks([0,1], ["预测正常", "预测异常"]); ax.set_yticks([0,1], ["实际正常", "实际异常"]); ax.set_title("混淆矩阵 · 可复现工程数据")
    for i in range(2):
        for j in range(2): ax.text(j, i, str(matrix[i][j]), ha="center", va="center", fontsize=18, color="white" if matrix[i][j] > max(map(max,matrix))/2 else "#10233c")
    fig.colorbar(image, ax=ax, fraction=.046); fig.tight_layout(); fig.savefig(EVAL_DIR / "confusion_matrix.png", dpi=180); plt.close(fig)

    history = training.get("history", [])
    if history:
        epochs = [int(row["epoch"]) for row in history]
        fig, ax1 = plt.subplots(figsize=(8, 4.6)); ax2 = ax1.twinx()
        ax1.plot(epochs, [row["train_loss"] for row in history], marker="o", color=colors["blue"], label="训练损失")
        ax2.plot(epochs, [row["validation_f1"] for row in history], marker="s", color=colors["green"], label="验证集 F1")
        ax1.set(xlabel="训练轮次", ylabel="损失", title="模型训练过程"); ax2.set_ylabel("F1 值"); ax2.set_ylim(0,1.05)
        lines = ax1.lines + ax2.lines; ax1.legend(lines, [line.get_label() for line in lines], loc="center right")
        fig.tight_layout(); fig.savefig(EVAL_DIR / "training_history.png", dpi=180); plt.close(fig)


def replacements(metrics, training):
    return {
        "{{DETECTION_RATE}}": f"{metrics['detection_rate']*100:.2f}%",
        "{{FPR}}": f"{metrics['false_positive_rate']*100:.2f}%",
        "{{PRECISION}}": f"{metrics['precision']*100:.2f}%",
        "{{F1}}": f"{metrics['f1']:.4f}",
        "{{AUC}}": f"{metrics['roc_auc']:.4f}",
        "{{LATENCY_MS}}": f"{metrics['mean_latency_ms']:.3f}",
        "{{THROUGHPUT}}": f"{metrics['throughput_sequences_per_second']:.1f}",
        "{{MODEL_MB}}": f"{metrics['checkpoint_size_mb']:.3f}",
        "{{SEQUENCE_COUNT}}": str(metrics["sequence_count"]),
        "{{PARAMETERS}}": f"{training['model_parameters']:,}",
    }


def set_east_asia_font(run, name: str, size: float | None = None, bold: bool | None = None):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.bold = bold


def configure_document(doc: Document, title: str):
    section = doc.sections[0]; section.page_height = Cm(29.7); section.page_width = Cm(21.0); section.top_margin = Cm(2.4); section.bottom_margin = Cm(2.4); section.left_margin = Cm(2.6); section.right_margin = Cm(2.4)
    normal = doc.styles["Normal"]; normal.font.name = "宋体"; normal._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体"); normal.font.size = Pt(12); normal.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE; normal.paragraph_format.first_line_indent = Cm(0.74); normal.paragraph_format.space_after = Pt(0)
    for name, size in [("Title",22),("Heading 1",16),("Heading 2",14),("Heading 3",12)]:
        style = doc.styles[name]; style.font.name = "黑体"; style._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体"); style.font.size = Pt(size); style.font.bold = True; style.font.color.rgb = RGBColor(0,0,0)
    header = section.header.paragraphs[0]; header.alignment = WD_ALIGN_PARAGRAPH.CENTER; run = header.add_run("2026第二届大学生人工智能安全竞赛 · 作品报告"); set_east_asia_font(run,"宋体",9)
    footer = section.footer.paragraphs[0]; footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer.add_run("第 "); set_east_asia_font(run,"宋体",9)
    fld = OxmlElement("w:fldSimple"); fld.set(qn("w:instr"), "PAGE"); footer._p.append(fld)
    run = footer.add_run(" 页"); set_east_asia_font(run,"宋体",9)
    doc.core_properties.title = title; doc.core_properties.subject = "2026第二届大学生人工智能安全竞赛自主命题赛道作品"
    doc.core_properties.author = ""; doc.core_properties.last_modified_by = ""; doc.core_properties.comments = ""
    update_fields = OxmlElement("w:updateFields")
    update_fields.set(qn("w:val"), "true")
    doc.settings._element.append(update_fields)


def add_cover(doc: Document):
    for _ in range(3): doc.add_paragraph()
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; r = p.add_run("2026第二届大学生人工智能安全竞赛"); set_east_asia_font(r,"黑体",22,True)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; r = p.add_run("作 品 报 告"); set_east_asia_font(r,"黑体",28,True)
    doc.add_paragraph(); p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; r = p.add_run("AgentGuard：基于行为序列建模的 AI Agent 异常检测系统"); set_east_asia_font(r,"黑体",18,True)
    doc.add_paragraph(); table = doc.add_table(rows=4, cols=2); table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = "Table Grid"
    values = [("作品类型","开放式自由命题"),("电子邮箱","【提交前填写】"),("提交日期","2026年【  】月【  】日"),("匿名提示","本报告不含学校、院系和指导教师信息")]
    for row, values_row in zip(table.rows, values):
        for cell, text in zip(row.cells, values_row): cell.text = text; cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    doc.add_page_break()
    p=doc.add_paragraph(); r=p.add_run("填写说明"); set_east_asia_font(r,"黑体",16,True); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    notes=["本作品为可运行、可演示、可验证的完整设计。","正文采用A4、宋体小四、1.5倍行距。","报告中的自建工程基准不代表生产环境效果。","初赛版本不得出现学校、院系、指导教师等身份信息。"]
    for note in notes: doc.add_paragraph(note, style="List Bullet")
    doc.add_page_break(); p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run("目  录"); set_east_asia_font(r,"黑体",16,True)
    toc = OxmlElement("w:fldSimple"); toc.set(qn("w:instr"), 'TOC \\o "1-3" \\h \\z \\u'); p=doc.add_paragraph(); p._p.append(toc); doc.add_page_break()


def add_table(doc: Document, rows: list[list[str]]):
    width = max(len(row) for row in rows); table = doc.add_table(rows=len(rows), cols=width); table.style = "Table Grid"; table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell=table.cell(i,j); cell.text=value.strip(); cell.vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.first_line_indent = Cm(0); paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if i==0 else WD_ALIGN_PARAGRAPH.LEFT
                for run in paragraph.runs: set_east_asia_font(run,"宋体",10.5,i==0)
    doc.add_paragraph()


def markdown_to_docx(md_path: Path, output_path: Path, token_values: dict[str,str]):
    text = md_path.read_text(encoding="utf-8")
    for key,value in token_values.items(): text=text.replace(key,value)
    lines=text.splitlines(); doc=Document(); configure_document(doc,"AgentGuard作品报告"); add_cover(doc)
    i=0; figure_number=0
    while i < len(lines):
        line=lines[i].strip()
        if not line: i+=1; continue
        if line.startswith("|") and i+1<len(lines) and re.match(r"^\|?\s*:?-+", lines[i+1].strip()):
            rows=[]; i+=2
            header=[cell.strip() for cell in line.strip("|").split("|")]; rows.append(header)
            while i<len(lines) and lines[i].strip().startswith("|"):
                rows.append([cell.strip() for cell in lines[i].strip().strip("|").split("|")]); i+=1
            add_table(doc,rows); continue
        image_match=re.match(r"!\[(.+?)\]\((.+?)\)",line)
        if image_match:
            path=ROOT/image_match.group(2); figure_number+=1
            if path.exists():
                p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run().add_picture(str(path),width=Inches(6.2))
                cap=doc.add_paragraph(f"图 {figure_number}  {image_match.group(1)}"); cap.alignment=WD_ALIGN_PARAGRAPH.CENTER; cap.paragraph_format.first_line_indent=Cm(0)
            i+=1; continue
        if line.startswith("# "):
            if doc.paragraphs and line[2:] not in {"摘要"}: doc.add_page_break()
            p=doc.add_heading(line[2:],level=1); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; i+=1; continue
        if line.startswith("## "): doc.add_heading(line[3:],level=2); i+=1; continue
        if line.startswith("### "): doc.add_heading(line[4:],level=3); i+=1; continue
        if re.match(r"^[-*] ",line):
            p=doc.add_paragraph(line[2:],style="List Bullet"); p.paragraph_format.first_line_indent=Cm(0); i+=1; continue
        if re.match(r"^\d+\. ",line):
            p=doc.add_paragraph(re.sub(r"^\d+\. ","",line),style="List Number"); p.paragraph_format.first_line_indent=Cm(0); i+=1; continue
        p=doc.add_paragraph(line.strip("`")); p.alignment=WD_ALIGN_PARAGRAPH.JUSTIFY; i+=1
    output_path.parent.mkdir(parents=True,exist_ok=True); doc.save(output_path)


def add_doc_title(doc, title, subtitle=None):
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run(title); set_east_asia_font(r,"黑体",20,True)
    if subtitle:
        p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run(subtitle); set_east_asia_font(r,"宋体",11)
    doc.add_paragraph()


def make_analysis_report(metrics, alerts, output_path):
    doc=Document(); configure_document(doc,"AgentGuard检测分析报告"); add_doc_title(doc,"AgentGuard 检测分析报告","自主命题赛道 · 自建增强工程基准")
    doc.add_heading("一、数据披露",1); doc.add_paragraph("本报告使用固定随机种子2026生成的可复现工程基准，目的是验证代码、阈值、解释和性能测量链路。自主命题赛道没有为本题提供数据，该结果不代表真实部署效果。")
    doc.add_heading("二、核心指标",1)
    rows=[["指标","数值"],["测试序列",str(metrics['sequence_count'])],["检出率",f"{metrics['detection_rate']*100:.2f}%"],["误报率",f"{metrics['false_positive_rate']*100:.2f}%"],["Precision",f"{metrics['precision']*100:.2f}%"],["F1",f"{metrics['f1']:.4f}"],["ROC-AUC",f"{metrics['roc_auc']:.4f}"],["平均延迟",f"{metrics['mean_latency_ms']:.3f} ms/序列"],["吞吐率",f"{metrics['throughput_sequences_per_second']:.1f} 序列/秒"],["模型大小",f"{metrics['checkpoint_size_mb']:.3f} MB"]]; add_table(doc,rows)
    for image,caption in [("confusion_matrix.png","混淆矩阵"),("score_distribution.png","异常分数分布"),("scenario_detection_rate.png","分场景检出覆盖")]:
        p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run().add_picture(str(EVAL_DIR/image),width=Inches(5.8)); cap=doc.add_paragraph(caption); cap.alignment=WD_ALIGN_PARAGRAPH.CENTER; cap.paragraph_format.first_line_indent=Cm(0)
    doc.add_heading("三、典型告警证据",1)
    for alert in alerts[:3]:
        doc.add_heading(f"{alert['entity_id']} · {alert['severity']} · {alert['score']}",2); doc.add_paragraph(alert['explanation']);
        evidence=[["时间","事件","对象","贡献"]]+[[e['timestamp'],e['event'],e['object_name'],str(e['contribution'])] for e in alert['evidence'][:4]]; add_table(doc,evidence)
    doc.add_heading("四、真实环境验证要求",1); doc.add_paragraph("在授权环境采集并脱敏真实轨迹，按时间和实体隔离训练/验证/测试；只在验证集校准阈值；补充LSTM和消融；报告逐日漂移、留出场景、误报/漏报案例及真实设备资源占用。")
    output_path.parent.mkdir(parents=True,exist_ok=True); doc.save(output_path)


def make_guide_figures():
    """依据项目真实输入格式和评测告警生成可复现的操作示意图。"""
    EVAL_DIR.mkdir(parents=True, exist_ok=True)
    dashboard = EVAL_DIR / "dashboard.png"
    if not dashboard.is_file():
        raise FileNotFoundError(f"使用说明所需主界面截图缺失：{dashboard}")

    input_line = {
        "timestamp": "2026-07-22T10:00:03Z",
        "entity_id": "agent-demo-01",
        "source": "agent",
        "event_type": "tool_call",
        "action": "read_file",
        "object_name": "credentials.txt",
        "result": "success",
    }
    fig, ax = plt.subplots(figsize=(12, 3.2), facecolor="#0d1b2a")
    ax.set_facecolor("#0d1b2a"); ax.axis("off")
    ax.text(.03, .82, "input.jsonl · UTF-8 · 每行一个事件", color="#51d5d0", fontsize=14, weight="bold", transform=ax.transAxes)
    ax.text(.03, .42, json.dumps(input_line, ensure_ascii=False), color="#eaf3ff", fontsize=10, family="monospace", wrap=True, transform=ax.transAxes)
    ax.text(.03, .10, "必填：timestamp / entity_id / event_type / action", color="#90a9c4", fontsize=11, transform=ax.transAxes)
    fig.tight_layout(); fig.savefig(EVAL_DIR / "guide_input_example.png", dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor()); plt.close(fig)

    metrics, _, alerts = load_inputs()
    alert = alerts[0] if alerts else {"entity_id": "agent-demo-01", "severity": "high", "score": 0.97, "explanation": "检测到异常行为序列。", "evidence": []}
    fig, ax = plt.subplots(figsize=(11, 4.8), facecolor="#f4f7fb"); ax.axis("off")
    ax.text(.04, .88, "告警结果", fontsize=19, weight="bold", color="#142b46", transform=ax.transAxes)
    ax.text(.04, .70, f"实体  {alert.get('entity_id')}     等级  {alert.get('severity')}     风险分  {float(alert.get('score', 0)):.3f}", fontsize=13, color="#c43d4f", transform=ax.transAxes)
    ax.text(.04, .49, str(alert.get("explanation", "")), fontsize=11, color="#314b68", wrap=True, transform=ax.transAxes)
    ax.text(.04, .18, "输出文件：artifacts/evaluation/alerts.jsonl", fontsize=11, color="#56708d", transform=ax.transAxes)
    fig.tight_layout(); fig.savefig(EVAL_DIR / "guide_alert_result.png", dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor()); plt.close(fig)

    model_score = float(alert.get("model_score", metrics.get("mean_anomaly_score", 0.65)))
    rule_score = float(alert.get("rule_score", 1.0 if alert.get("score", 0) else 0.0))
    hybrid_score = float(alert.get("score", max(model_score, rule_score)))
    fig, axes = plt.subplots(1, 3, figsize=(11, 3.8), facecolor="#f4f7fb")
    for ax, title, value, color in zip(axes, ("模型分", "规则分", "混合分"), (model_score, rule_score, hybrid_score), ("#5d8ffc", "#f6bd60", "#e45b5b")):
        ax.barh([0], [value], color=color, height=.35); ax.set_xlim(0, 1); ax.set_yticks([]); ax.set_title(title, weight="bold"); ax.text(min(value + .02, .92), 0, f"{value:.3f}", va="center"); ax.grid(axis="x", alpha=.2)
    fig.suptitle("双证据评分与最终决策", fontsize=17, weight="bold", color="#142b46")
    fig.tight_layout(); fig.savefig(EVAL_DIR / "guide_score_evidence.png", dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor()); plt.close(fig)


def add_guide_figure(doc: Document, path: Path, caption: str, width: float = 6.15):
    paragraph = doc.add_paragraph(); paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run().add_picture(str(path), width=Inches(width))
    caption_paragraph = doc.add_paragraph(caption); caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    caption_paragraph.paragraph_format.first_line_indent = Cm(0)


def make_user_guide(output_path):
    make_guide_figures()
    doc=Document(); configure_document(doc,"AgentGuard系统使用说明"); add_doc_title(doc,"AgentGuard 系统使用说明","版本 1.0")
    doc.add_heading("1. 环境与一键构建",1)
    for item in ("Windows 10/11 或 macOS；Python 3.11/3.12；CPU 即可，建议 8GB 内存。", "安装：Windows 执行 setup_env.ps1；macOS 执行 setup_env.sh。", "完整复现：Windows 执行 run_all.ps1；macOS 执行 run_all.sh。"):
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("2. 系统启动与主界面",1)
    doc.add_paragraph("Windows 执行 run_demo.ps1、macOS 执行 run_demo.sh，浏览器访问 http://127.0.0.1:8080；页面加载模型、指标和告警数据后即可点击“运行演示检测”。")
    add_guide_figure(doc, EVAL_DIR / "dashboard.png", "图 1  系统启动与主界面")
    doc.add_heading("3. 日志输入示例",1)
    doc.add_paragraph("自有日志采用 UTF-8 JSONL。每个实体的事件按 timestamp 排序，适配层再切分为长度 24、步长 8 的行为窗口。")
    add_guide_figure(doc, EVAL_DIR / "guide_input_example.png", "图 2  日志输入示例")
    doc.add_paragraph("命令行分析：python scripts/analyze.py input.jsonl。输入仅限已授权且已脱敏的数据，不要写入 API Key、Token、Cookie 或个人身份信息。")
    doc.add_heading("4. 告警结果页面",1)
    doc.add_paragraph("告警同时保留实体、严重等级、风险分、解释和关键事件，可从 artifacts/evaluation/alerts.jsonl 回查原始证据。")
    add_guide_figure(doc, EVAL_DIR / "guide_alert_result.png", "图 3  告警结果页面")
    doc.add_heading("5. 模型分、规则分与证据链",1)
    doc.add_paragraph("模型分反映行为序列偏离，规则分反映透明的高风险顺序命中，混合分用于最终决策。97% 检出率属于混合系统，不是 Transformer 单模型成绩。")
    add_guide_figure(doc, EVAL_DIR / "guide_score_evidence.png", "图 4  模型分、规则分与证据链")
    doc.add_heading("6. 自有数据训练与故障处理",1)
    for item in ("把授权且脱敏的 train/validation/test.jsonl 放入 data/local，再运行 train.py 与 evaluate.py。", "端口占用：使用 serve.py --port 8081。", "模型缺失：先运行 generate_demo_data.py、train.py。", "字段错误：检查 timestamp、entity_id、event_type、action。", "现场断网不影响核心功能。"):
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("7. 部署安全边界",1)
    doc.add_paragraph("服务默认仅监听本机。真实部署前应增加认证、TLS、上传限制、日志脱敏、访问审计与模型漂移监测；当前自建基准结果不能直接外推到生产环境。")
    output_path.parent.mkdir(parents=True,exist_ok=True); doc.save(output_path); scrub_office_metadata(output_path)


def ppt_text(slide, text, x, y, w, h, size=20, color="DDEBFA", bold=False, align=PP_ALIGN.LEFT):
    box=slide.shapes.add_textbox(PptInches(x),PptInches(y),PptInches(w),PptInches(h)); frame=box.text_frame; frame.clear(); frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=frame.paragraphs[0]; p.alignment=align; run=p.add_run(); run.text=text; run.font.name="Microsoft YaHei"; run.font.size=PptPt(size); run.font.bold=bold; run.font.color.rgb=PptRGB.from_string(color); return box


def base_slide(prs, title, number):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=PptRGB(8,17,31)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,PptInches(.12),prs.slide_height).fill.solid(); slide.shapes[-1].fill.fore_color.rgb=PptRGB(81,213,208); slide.shapes[-1].line.fill.background()
    ppt_text(slide,title,.55,.25,11.6,.55,24,"EAF3FF",True); ppt_text(slide,f"{number:02d}",12.25,.28,.55,.4,11,"6F89A8",True,PP_ALIGN.RIGHT)
    ppt_text(slide,"AgentGuard · 2026第二届大学生人工智能安全竞赛",.55,7.12,8,.2,8,"57708D")
    return slide


def card(slide,x,y,w,h,title,body,accent="51D5D0"):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptInches(x),PptInches(y),PptInches(w),PptInches(h)); shape.fill.solid(); shape.fill.fore_color.rgb=PptRGB(18,32,55); shape.line.color.rgb=PptRGB.from_string("29415F")
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptInches(x),PptInches(y),PptInches(.06),PptInches(h)).fill.solid(); slide.shapes[-1].fill.fore_color.rgb=PptRGB.from_string(accent); slide.shapes[-1].line.fill.background()
    ppt_text(slide,title,x+.22,y+.16,w-.4,.35,15,"E8F2FD",True); ppt_text(slide,body,x+.22,y+.58,w-.4,h-.68,11,"91A8C2")


def make_ppt(metrics, training, alerts, output_path):
    prs=Presentation(); prs.slide_width=PptInches(13.333); prs.slide_height=PptInches(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=PptRGB(8,17,31)
    slide.shapes.add_shape(MSO_SHAPE.OVAL,PptInches(9.7),PptInches(-1.2),PptInches(5),PptInches(5)).fill.solid(); slide.shapes[-1].fill.fore_color.rgb=PptRGB(26,57,96); slide.shapes[-1].fill.transparency=25; slide.shapes[-1].line.fill.background()
    ppt_text(slide,"AgentGuard",.75,1.25,7,1,38,"51D5D0",True); ppt_text(slide,"基于行为序列建模的 AI Agent 异常检测系统",.78,2.25,10.7,.7,26,"EDF5FF",True); ppt_text(slide,"Transformer · 顺序规则 · 可解释证据 · CPU离线部署",.8,3.15,9,.45,15,"90A9C4"); ppt_text(slide,"自主命题赛道 · AI 衍生安全方向",.8,5.85,5,.35,13,"5D8FFC",True)
    s=base_slide(prs,"为什么单事件规则看不见 Agent 风险链",2); card(s,.7,1.15,3.7,2.05,"上下文缺失","读取文件、连接网络、写注册表单独看都可能合法，风险存在于动作组合与顺序。","FF6B78"); card(s,4.8,1.15,3.7,2.05,"行为漂移","用户、Agent与业务周期不断变化，静态阈值容易产生告警疲劳。","F6BD60"); card(s,8.9,1.15,3.7,2.05,"解释断层","只给一个风险分无法帮助运维判断；必须回到关键片段和原始日志。","51D5D0"); ppt_text(s,"不可信内容 → 权限提升 → 读取凭据 → 新外连 → 上传",1.1,4.25,11.1,.65,24,"EAF3FF",True,PP_ALIGN.CENTER); ppt_text(s,"把一句输入追踪为一串可核验的主机行为",2.2,5.15,9,.4,16,"8FA6C1",False,PP_ALIGN.CENTER)
    s=base_slide(prs,"自主命题作品的四类证据",3); specs=[("问题价值","可验证","真实风险链与应用边界"),("方法有效","可对比","模型 / 规则 / 混合分层"),("工程质量","可运行","CPU延迟 / 测试 / 一键复现"),("创新解释","可追溯","关键事件 + 原始日志 + 战术")]
    for i,(a,b,c) in enumerate(specs): card(s,.65+i*3.15,1.25,2.8,3.9,a,b+"\n\n"+c,["FF6B78","5D8FFC","F6BD60","51D5D0"][i]); ppt_text(s,"内部阈值目标：验证集 FPR≤15%；该数值不是赛道官方门槛",1.1,5.65,11.1,.5,18,"59D39B",True,PP_ALIGN.CENTER)
    s=base_slide(prs,"端到端架构：离线、可复现、证据保全",4); s.shapes.add_picture(str(EVAL_DIR/"model_architecture.png"),PptInches(.7),PptInches(1.15),width=PptInches(11.9)); ppt_text(s,"适配层隔离日志厂商格式；训练与推理共享同一词表、配置和阈值",1.1,6.15,11,.4,14,"8FA6C1",False,PP_ALIGN.CENTER)
    s=base_slide(prs,"统一 Agent + 主机行为序列",5); events=[("模型","接收输入"),("规划","推理决策"),("工具","读取凭据"),("主机","网络连接"),("工具","上传数据")]
    for i,(a,b) in enumerate(events):
        card(s,.65+i*2.5,1.45,2.1,1.55,a,b,"51D5D0" if i<2 else "FF6B78");
        if i<4: ppt_text(s,"→",2.78+i*2.5,1.87,.3,.4,21,"5D8FFC",True,PP_ALIGN.CENTER)
    card(s,1.0,4.15,3.3,1.55,"行为词元","source | type | action | object | result","5D8FFC"); card(s,5.0,4.15,3.3,1.55,"连续特征","时间间隔、昼夜周期、结果、对象长度","F6BD60"); card(s,9.0,4.15,3.3,1.55,"防泄漏","label / scenario / risk_hint 不进入模型","59D39B")
    s=base_slide(prs,"轻量 Transformer：分类 + 下一事件预测",6); s.shapes.add_picture(str(EVAL_DIR/"model_architecture.png"),PptInches(.8),PptInches(1.1),width=PptInches(7.4)); card(s,8.65,1.35,3.8,1.45,"模型配置",f"2层 · 4头 · 隐藏维64\n参数 {training['model_parameters']:,} · {metrics['checkpoint_size_mb']:.3f} MB","51D5D0"); card(s,8.65,3.15,3.8,1.45,"多任务损失","加权分类损失 + 0.25 × 下一事件损失","5D8FFC"); card(s,8.65,4.95,3.8,1.2,"异常分数","0.8 × 分类概率 + 0.2 × 行为惊异度","F6BD60")
    s=base_slide(prs,"解释不是一张热力图，而是一条证据链",7); sample=alerts[0] if alerts else {"explanation":"关键行为序列显著偏离正常上下文。","evidence":[],"mapped_tactics":["Behavioral Anomaly"]}; card(s,.7,1.15,5.3,4.9,"自然语言研判",sample['explanation']+"\n\n战术："+" / ".join(sample.get('mapped_tactics',[])),"FF6B78"); evidence=sample.get('evidence',[])[:4]
    for i,e in enumerate(evidence): card(s,6.35,1.15+i*1.22,6.25,1.0,f"#{i+1}  {e['event']}",f"{e['timestamp']} · {e['object_name']} · contribution={e['contribution']}","51D5D0" if i==0 else "5D8FFC")
    s=base_slide(prs,"可复现实验：分层披露，不把混合结果归给模型",8); s.shapes.add_picture(str(EVAL_DIR/"score_distribution.png"),PptInches(.65),PptInches(1.1),width=PptInches(6.6)); metrics_text=[("Recall",f"{metrics['detection_rate']*100:.1f}%"),("FPR",f"{metrics['false_positive_rate']*100:.1f}%"),("F1",f"{metrics['f1']:.3f}"),("Latency",f"{metrics['mean_latency_ms']:.3f} ms")]
    for i,(a,b) in enumerate(metrics_text): card(s,7.55+(i%2)*2.55,1.35+(i//2)*1.85,2.25,1.45,a,b,["59D39B","FF6B78","5D8FFC","F6BD60"][i]); ppt_text(s,"披露：固定种子自建工程基准，不代表真实部署效果",7.5,5.35,5.2,.7,13,"FFB2BA",True,PP_ALIGN.CENTER)
    s=base_slide(prs,"现场演示：从日志到可追溯告警",9)
    dashboard=EVAL_DIR/"dashboard.png"
    if dashboard.exists(): s.shapes.add_picture(str(dashboard),PptInches(.55),PptInches(1.0),width=PptInches(8.55))
    card(s,9.35,1.15,3.35,1.35,"1  一键启动","run_demo.ps1 / .sh 离线加载","51D5D0"); card(s,9.35,2.85,3.35,1.35,"2  运行检测","展示风险分、阈值与场景","5D8FFC"); card(s,9.35,4.55,3.35,1.35,"3  核验证据","查看关键事件与原始日志","F6BD60")
    s=base_slide(prs,"四个可验证创新点",10); innovations=[("双域序列","Agent工具调用与主机遥测统一"),("双任务双分数","分类与正常语法偏离互补"),("FPR约束校准","直接优化官方高精度目标"),("证据闭环","注意力+惊异度+原始日志")]
    for i,(a,b) in enumerate(innovations): card(s,.8+(i%2)*6.1,1.25+(i//2)*2.45,5.55,1.95,a,b,["51D5D0","5D8FFC","F6BD60","FF6B78"][i])
    s=base_slide(prs,"已知边界与真实验证路线",11); card(s,.8,1.2,5.6,4.8,"我们主动披露","• 当前结果来自自建数据\n• 工程集模板边界清晰\n• 注意力不是因果解释\n• 大规模日志需实体回收策略","FF6B78"); card(s,6.95,1.2,5.6,4.8,"授权轨迹接入后","• 时间 + 实体隔离\n• LSTM / Transformer / 混合消融\n• 多个留一场景测试\n• 逐日漂移 + 真实资源测试","59D39B")
    s=base_slide(prs,"AgentGuard：把 AI 输入追踪为可核验的行为证据",12); ppt_text(s,"可运行",1.1,1.45,3,.6,28,"51D5D0",True,PP_ALIGN.CENTER); ppt_text(s,"可检测",5.15,1.45,3,.6,28,"5D8FFC",True,PP_ALIGN.CENTER); ppt_text(s,"可解释",9.2,1.45,3,.6,28,"F6BD60",True,PP_ALIGN.CENTER); ppt_text(s,"源码 · 模型 · 测试 · 告警 · Web · 文档\n一键复现，离线演示，原始日志可追溯",1.35,3.15,10.7,1.2,25,"EDF5FF",True,PP_ALIGN.CENTER); ppt_text(s,"谢谢 · 请专家指导",3.7,5.55,6,.55,20,"8FA6C1",False,PP_ALIGN.CENTER)
    output_path.parent.mkdir(parents=True,exist_ok=True); prs.save(output_path)


def fill_originality(output_path):
    source=next((ROOT/"00_比赛原始材料").glob("originality-declaration.docx")); doc=Document(source)
    for paragraph in doc.paragraphs:
        if "作品报告" in paragraph.text and "郑重声明" in paragraph.text:
            paragraph.text=paragraph.text.replace("作品报告                                                                ，","作品报告《AgentGuard：基于行为序列建模的AI Agent异常检测系统》，")
            for run in paragraph.runs: set_east_asia_font(run,"宋体",12)
    doc.core_properties.author = ""; doc.core_properties.last_modified_by = ""; doc.core_properties.comments = ""
    doc.save(output_path); scrub_office_metadata(output_path)


def make_enhanced_copy(source_path: Path, output_path: Path):
    doc = Document(source_path)
    accent = RGBColor(31, 78, 121)
    for style_name in ["Title", "Heading 1", "Heading 2", "Heading 3"]:
        doc.styles[style_name].font.color.rgb = accent
    for table in doc.tables:
        if not table.rows:
            continue
        for cell in table.rows[0].cells:
            shading = OxmlElement("w:shd")
            shading.set(qn("w:fill"), "D9EAF7")
            cell._tc.get_or_add_tcPr().append(shading)
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
                    run.font.color.rgb = RGBColor(24, 55, 86)
    doc.core_properties.title = "AgentGuard作品报告（评审增强排版版）"
    doc.save(output_path)


def make_manifest(output_path):
    files=[]
    for path in sorted(ROOT.rglob("*")):
        if (
            path.is_file()
            and not any(
                part in {".git", ".venv", "__pycache__", ".matplotlib", "archive_pre_optimization", "superpowers"}
                for part in path.parts
            )
            and not path.name.endswith((".pyc", ".pyo"))
            and path.name != "torch-2.7.1+cpu-cp311-cp311-win_amd64.whl"
        ):
            files.append(f"{path.relative_to(ROOT)}\t{path.stat().st_size} bytes")
    output_path.write_text("AgentGuard source and artifact manifest\n\n"+"\n".join(files),encoding="utf-8")


def main():
    metrics,training,alerts=load_inputs(); make_figures(metrics,training); values=replacements(metrics,training); SUBMISSION_DIR.mkdir(parents=True,exist_ok=True)
    official_report = SUBMISSION_DIR/"作品报告_AgentGuard_自主命题模板版.docx"
    build_report(official_report, values)
    make_user_guide(SUBMISSION_DIR/"系统使用说明_AgentGuard.docx")
    fill_originality(SUBMISSION_DIR/"作品原创性声明_待签字盖章.docx")
    final_ppt = SUBMISSION_DIR/"答辩PPT_AgentGuard_自主命题最终版.pptx"
    if final_ppt.is_file():
        scrub_office_metadata(final_ppt)
    shutil.copy2(ROOT/"docs"/"06_答辩演示稿与问答.md",SUBMISSION_DIR/"答辩演示稿与问答.md")
    shutil.copy2(ROOT/"docs"/"07_最终提交清单.md",SUBMISSION_DIR/"最终提交清单.md")
    make_manifest(SUBMISSION_DIR/"项目文件清单.txt")
    summary={
        "generated_files":[path.name for path in SUBMISSION_DIR.iterdir() if path.is_file()],
        "data_disclosure":metrics["dataset_disclosure"],
        "track":"自主命题赛道",
        "organizer_dataset_provided":False,
        "real_environment_validation_recommended":True,
        "ppt_note":"最终PPT由独立演示文稿工具生成，本脚本不会覆盖。",
    }
    (SUBMISSION_DIR/"材料生成说明.json").write_text(json.dumps(summary,ensure_ascii=False,indent=2),encoding="utf-8")
    print(json.dumps(summary,ensure_ascii=False,indent=2))


if __name__ == "__main__":
    main()
